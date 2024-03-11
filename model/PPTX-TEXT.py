from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob
import os
import re
n = 0
l2=[]
def write_image(shape, fname):
    global n
    image = shape.image
    image_bytes = image.blob
    image_filename = fname[:-5] + '{:03d}.{}'.format(n, image.ext)
    n += 1
    print(image_filename)
    os.chdir("/readpptx/images")
    with open(image_filename, 'wb') as f:
        f.write(image_bytes)
    os.chdir("/readpptx")
def visitor(shape, fname, file):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        print(" ")
    elif hasattr(shape, "text"):
        try:
            text = shape.text
            #print(text)
            a=re.sub(r'[#]','',text)
            print(a)
            b=a.split()
            l2.append(len(b))
            file.write(text + "\n")
        except Exception as e:
            print(f"Error writing text: {e}")
def iter_shapes(prs, fname, file):
    for slide in prs.slides:
        for shape in slide.shapes:
            visitor(shape, fname, file)
file = open("MyFile.txt", "a+", encoding="utf-8")
for each_file in glob.glob("*.pptx"):
    fname = os.path.basename(each_file)
    #file.write("-------------------" + fname + "----------------------\n")
    prs = Presentation(each_file)
    #print("---------------" + fname + "-------------------")
    iter_shapes(prs, fname, file)
file.close()
print("Total Tokens",sum(l2))