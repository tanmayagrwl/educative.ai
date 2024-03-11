import re
from typing import BinaryIO, TypedDict
import google.generativeai as ai
from model.genai import safety_settings
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from transformers import pipeline

summarizer = pipeline("summarization", model="Falconsai/text_summarization")

generation_config = ai.GenerationConfig(
    temperature=0.4,
    top_p=1,
    top_k=32,
    max_output_tokens=2000,
)

model = ai.GenerativeModel(
    model_name="gemini-pro-vision",
    generation_config=generation_config,
    safety_settings=safety_settings,
)


class ImageData(TypedDict):
    mime_type: str
    data: bytes


def get_ocr_text(file: ImageData):
    return model.generate_content(
        [
            "What is written on this image?",
            {
                "mime_type": file["mime_type"],
                "data": file["data"],
            },
        ]
    ).text


class PPTData(TypedDict):
    file: BinaryIO


def ocr_ppt(file: PPTData):
    l2=[]
    l3=[]
    def visitor(shape, fname, file):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(" ")
        elif hasattr(shape, "text"):
            try:
                text = shape.text
                #print(text)
                a=re.sub(r'[#]','',text)
                l3.append(a)
                b=a.split()
                l2.append(len(b))
                file.write(text + "\n")
            except Exception as e:
                print(f"Error writing text: {e}")

    def iter_shapes(prs, fname, file):
        for slide in prs.slides:
            for shape in slide.shapes:
                visitor(shape, fname, file)

    prs = Presentation(file["file"])

    iter_shapes(prs, file["file"].name or "Unknown", file["file"])
    rw=' '.join(l3)

    return re.sub(r'\s+', ' ', rw).strip()

def sum_ppt(content: str):
    sum = summarizer(content, max_length=1000, min_length=30, do_sample=False)[0]["summary_text"]
    return re.sub(r'\s+', ' ', sum).strip()